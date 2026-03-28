from __future__ import annotations

from pathlib import Path
from typing import List

from docx import Document
from pptx import Presentation

from autograde.extractors.base import BaseExtractor
from autograde.models import Artifact, EvidenceObject, Location


class DocxTextExtractor(BaseExtractor):
    extractor_id = "docx_text_extractor_v1"
    supported_types = {"docx_report"}

    def extract(self, artifact: Artifact) -> List[EvidenceObject]:
        doc = Document(artifact.storage_path)
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        evidence: List[EvidenceObject] = []
        file_name = Path(artifact.storage_path).name
        for idx, para in enumerate(paras, start=1):
            evidence.append(
                EvidenceObject(
                    evidence_id=f"{artifact.artifact_id}_docx_{idx}",
                    submission_id=artifact.submission_id,
                    artifact_id=artifact.artifact_id,
                    modality="text",
                    subtype="paragraph",
                    content=para,
                    structured_content={"tokens": len(para.split()), "source_format": "docx"},
                    preview=para[:120],
                    location=Location(file=file_name, paragraph=idx),
                    confidence=0.99,
                    extractor_id=self.extractor_id,
                    tags=["text_chunk", "docx"],
                )
            )
        evidence.append(
            EvidenceObject(
                evidence_id=f"{artifact.artifact_id}_docx_meta_1",
                submission_id=artifact.submission_id,
                artifact_id=artifact.artifact_id,
                modality="metadata",
                subtype="document_structure",
                content=None,
                structured_content={"paragraph_count": len(paras), "source_format": "docx"},
                preview=file_name,
                location=Location(file=file_name),
                confidence=0.98,
                extractor_id=self.extractor_id,
                tags=["document", "docx"],
            )
        )
        return evidence


class SlideDeckExtractor(BaseExtractor):
    extractor_id = "slide_deck_extractor_v1"
    supported_types = {"slide_deck"}

    def extract(self, artifact: Artifact) -> List[EvidenceObject]:
        prs = Presentation(artifact.storage_path)
        evidence: List[EvidenceObject] = []
        titles = []
        file_name = Path(artifact.storage_path).name
        for idx, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    texts.append(shape.text.strip())
            slide_text = "\n".join(texts).strip()
            title = texts[0] if texts else f"Slide {idx}"
            titles.append(title)
            if slide_text:
                evidence.append(
                    EvidenceObject(
                        evidence_id=f"{artifact.artifact_id}_slide_{idx}",
                        submission_id=artifact.submission_id,
                        artifact_id=artifact.artifact_id,
                        modality="text",
                        subtype="slide_text",
                        content=slide_text,
                        structured_content={"slide_index": idx, "shape_text_count": len(texts)},
                        preview=title[:120],
                        location=Location(file=file_name, page=idx),
                        confidence=0.98,
                        extractor_id=self.extractor_id,
                        tags=["slides", "presentation"],
                    )
                )
        evidence.append(
            EvidenceObject(
                evidence_id=f"{artifact.artifact_id}_slides_meta_1",
                submission_id=artifact.submission_id,
                artifact_id=artifact.artifact_id,
                modality="metadata",
                subtype="presentation_structure",
                content=None,
                structured_content={"slide_count": len(prs.slides), "titles": titles[:20]},
                preview=file_name,
                location=Location(file=file_name),
                confidence=0.98,
                extractor_id=self.extractor_id,
                tags=["slides", "presentation", "structure"],
            )
        )
        return evidence
