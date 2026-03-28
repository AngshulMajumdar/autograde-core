from __future__ import annotations

import json
from pathlib import Path


def build_sample_submission(root: str) -> None:
    base = Path(root)
    base.mkdir(parents=True, exist_ok=True)
    (base / "report.txt").write_text(
        """Title: Short Technical Report

We implemented Dijkstra's algorithm because it gives the shortest path in graphs with nonnegative weights. Therefore the approach is suitable here.

The code uses an adjacency list and a priority queue. However, the report does not discuss all edge cases in detail.

We conclude that the method works well on the provided inputs.""",
        encoding="utf-8",
    )
    (base / "main.py").write_text(
        """import heapq

def dijkstra(graph, src):
    dist = {node: float('inf') for node in graph}
    dist[src] = 0
    pq = [(0, src)]
    while pq:
        d, u = heapq.heappop(pq)
        if d > dist[u]:
            continue
        for v, w in graph[u]:
            nd = d + w
            if nd < dist[v]:
                dist[v] = nd
                heapq.heappush(pq, (nd, v))
    return dist
""",
        encoding="utf-8",
    )
    notebook = {
        "cells": [
            {"cell_type": "markdown", "source": ["# Experiment notes\n", "The algorithm was tested on a toy graph."]},
            {"cell_type": "code", "source": ["print('demo')\n"]},
        ],
        "metadata": {},
        "nbformat": 4,
        "nbformat_minor": 5,
    }
    (base / "analysis.ipynb").write_text(json.dumps(notebook), encoding="utf-8")
    (base / "results.csv").write_text("dataset,accuracy\ntoy,0.95\n", encoding="utf-8")


def build_contradictory_submission(root: str) -> None:
    base = Path(root)
    base.mkdir(parents=True, exist_ok=True)
    (base / "report.txt").write_text(
        """Title: Contradictory Technical Report

We implemented BFS algorithm for shortest paths. The accuracy was 98% on the benchmark dataset.

The method is simple and the algorithm works well in all cases.
""",
        encoding="utf-8",
    )
    (base / "main.py").write_text(
        """import heapq

def dijkstra(graph, src):
    dist = {node: float('inf') for node in graph}
    dist[src] = 0
    pq = [(0, src)]
    while pq:
        d, u = heapq.heappop(pq)
        if d > dist[u]:
            continue
        for v, w in graph[u]:
            nd = d + w
            if nd < dist[v]:
                dist[v] = nd
                heapq.heappush(pq, (nd, v))
    return dist
""",
        encoding="utf-8",
    )
    notebook = {
        "cells": [
            {"cell_type": "markdown", "source": ["# Notes\n", "The code was tested briefly.\n"]},
            {"cell_type": "code", "source": ["print('demo')\n"]},
        ],
        "metadata": {},
        "nbformat": 4,
        "nbformat_minor": 5,
    }
    (base / "analysis.ipynb").write_text(json.dumps(notebook), encoding="utf-8")
    (base / "results.csv").write_text("dataset,accuracy\nbench,0.55\n", encoding="utf-8")


def build_hardcoded_but_probeable_submission(root: str) -> None:
    base = Path(root)
    base.mkdir(parents=True, exist_ok=True)
    (base / "report.txt").write_text(
        """Title: Probeable Programming Report

We implemented Dijkstra's algorithm for shortest paths. The implementation uses a priority queue and should return correct path distances on small graphs.

The top-level script still contains machine-specific file paths, so direct execution may fail, but the core function is present.
""",
        encoding="utf-8",
    )
    (base / "main.py").write_text(
        """import heapq

DATA_PATH = '/Users/student/Desktop/private/input.csv'


def run_full_pipeline():
    with open(DATA_PATH, 'r', encoding='utf-8') as f:
        return f.read()


def dijkstra(graph, src):
    dist = {node: float('inf') for node in graph}
    dist[src] = 0
    pq = [(0, src)]
    while pq:
        d, u = heapq.heappop(pq)
        if d > dist[u]:
            continue
        for v, w in graph[u]:
            nd = d + w
            if nd < dist[v]:
                dist[v] = nd
                heapq.heappush(pq, (nd, v))
    return dist


if __name__ == '__main__':
    print(run_full_pipeline())
""",
        encoding="utf-8",
    )
    notebook = {
        "cells": [
            {"cell_type": "markdown", "source": ["# Notes\n", "Unit probes should still recover the callable function.\n"]},
        ],
        "metadata": {},
        "nbformat": 4,
        "nbformat_minor": 5,
    }
    (base / "analysis.ipynb").write_text(json.dumps(notebook), encoding="utf-8")
    (base / "results.csv").write_text("dataset,accuracy\ntoy,0.95\n", encoding="utf-8")


def build_humanities_short_answer_submission(root: str) -> None:
    base = Path(root)
    base.mkdir(parents=True, exist_ok=True)
    (base / "response.txt").write_text(
        """Response:

The passage argues that institutions matter because they shape how power is exercised. For example, a formal rule may look neutral, but in practice it can favor one group. However, the author also suggests that institutions are not enough by themselves, because social norms influence how rules are interpreted.""",
        encoding="utf-8",
    )



def build_humanities_essay_submission(root: str) -> None:
    base = Path(root)
    base.mkdir(parents=True, exist_ok=True)
    (base / "essay.txt").write_text(
        """Title: Interpreting the Passage

This essay argues that the passage presents institutions not as static structures but as contested arrangements whose effects depend on interpretation. The central claim is that rules matter, but they matter through lived practice.

According to the author, formal institutions distribute power by shaping what actions appear legitimate. This suggests that the text is not only describing rules; it is interpreting how authority becomes naturalized. For example, the passage shows that a rule can look universal while still reproducing an unequal outcome.

However, the text also resists a purely institutional explanation. The author indicates that norms, habits, and social expectations mediate how formal structures operate. This qualification matters because it prevents an overly mechanical reading of the passage.

A strong reading of the passage therefore requires both evidence and interpretation. The quoted examples support the idea that institutions guide conduct, while the broader discussion of social norms shows why the same rule may produce different consequences across settings.

In conclusion, the passage argues that institutions matter, but only through the social practices that animate them. That interpretation is more convincing than a narrow legal reading because it explains both stability and variation in outcomes.

References
Author, A. (2020). Interpreting Institutions.
""",
        encoding="utf-8",
    )



def build_engineering_circuit_submission(root: str) -> None:
    base = Path(root)
    base.mkdir(parents=True, exist_ok=True)
    (base / "report.txt").write_text(
        """Title: Active Filter Design Report

This submission uses an op-amp feedback design to realize the required low-pass behavior. The input signal enters the op-amp stage, and the RC feedback path sets the response.

The design explanation states that the topology is one valid active-filter family rather than a unique diagram. The target behavior is low-pass filtering with stable output gain in the pass band.

Simulation results are reported through gain and cutoff values, and the report argues that those values are consistent with the intended behavior.
""",
        encoding="utf-8",
    )
    (base / "circuit.svg").write_text(
        """<svg xmlns='http://www.w3.org/2000/svg' width='640' height='320'>
  <g id='input_node'><text x='20' y='80'>Input</text></g>
  <g id='resistor_R1'><text x='120' y='80'>R1</text></g>
  <g id='capacitor_C1'><text x='220' y='140'>C1</text></g>
  <g id='op_amp_stage'><text x='300' y='100'>OP_AMP</text></g>
  <g id='feedback_path'><text x='260' y='60'>feedback</text></g>
  <g id='output_node'><text x='520' y='100'>Output</text></g>
  <g id='ground'><text x='200' y='210'>ground</text></g>
  <text x='340' y='35'>active lowpass family</text>
</svg>
""",
        encoding="utf-8",
    )
    (base / "simulation.csv").write_text("""frequency,gain,cutoff
100,0.98,1000
1000,0.71,1000
5000,0.18,1000
""", encoding="utf-8")


def build_engineering_plausible_unknown_submission(root: str) -> None:
    base = Path(root)
    base.mkdir(parents=True, exist_ok=True)
    (base / "report.txt").write_text(
        """Title: Alternative Circuit Design

The circuit uses a nonstandard arrangement but still includes an op-amp, input path, output node, and a feedback relation. The report explains the design intention and the target behavior.
""",
        encoding="utf-8",
    )
    (base / "circuit.svg").write_text(
        """<svg xmlns='http://www.w3.org/2000/svg' width='640' height='320'>
  <g id='input_node'><text x='20' y='80'>Input</text></g>
  <g id='resistor_Rx'><text x='120' y='80'>Rx</text></g>
  <g id='op_amp_alt'><text x='300' y='100'>OPAMP</text></g>
  <g id='feedback_path'><text x='260' y='60'>feedback</text></g>
  <g id='output_node'><text x='520' y='100'>Output</text></g>
</svg>
""",
        encoding="utf-8",
    )
    (base / "simulation.csv").write_text("""frequency,gain
100,0.94
1000,0.76
""", encoding="utf-8")


def build_lab_science_submission(root: str) -> None:
    base = Path(root)
    base.mkdir(parents=True, exist_ok=True)
    (base / "report.txt").write_text(
        """Title: Enzyme Activity Lab Report

Method: We measured enzyme activity at different temperatures using the same assay volume and reaction time. The setup section states the temperature conditions and the observed reaction-rate variable.

Results: The observations show that activity increases from the low-temperature condition to a middle condition and then drops at the highest temperature. Therefore the data suggest an optimal intermediate temperature.

Interpretation: This suggests that both insufficient thermal energy and denaturation can reduce activity. A limitation is that measurement uncertainty and timing error may affect the precise optimum.
""",
        encoding="utf-8",
    )
    (base / "results.csv").write_text(
        "temperature,rate\n20,0.42\n30,0.63\n40,0.77\n50,0.51\n",
        encoding="utf-8",
    )


def build_mathematics_proof_submission(root: str) -> None:
    base = Path(root)
    base.mkdir(parents=True, exist_ok=True)
    (base / "solution.txt").write_text(
        """Title: Proof Submission

We prove that the sum of two even integers is even. Let a = 2m and b = 2n for integers m and n. Since a and b are even, they each have the stated form.

Therefore,
 a + b = 2m + 2n = 2(m+n).
Because m+n is an integer, the quantity a+b is divisible by 2. Hence the sum is even, which proves the claim.
""",
        encoding="utf-8",
    )



def build_multimodal_submission(root: str) -> None:
    base = Path(root)
    base.mkdir(parents=True, exist_ok=True)
    # docx report
    from docx import Document
    doc = Document()
    doc.add_heading('Project Reflection', level=1)
    doc.add_paragraph('This document summarizes the experiment and explains the observed improvement in accuracy.')
    doc.add_paragraph('The presentation and transcript provide additional context for the oral submission.')
    doc.save(str(base / 'report.docx'))

    # pptx deck
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = 'Results Summary'
    slide.placeholders[1].text = 'Accuracy improved from 0.81 to 0.89.\nKey limitation: small dataset.'
    prs.save(str(base / 'deck.pptx'))

    # wav audio with transcript
    import wave, struct, math
    wav_path = base / 'oral.wav'
    framerate = 8000
    duration = 0.5
    with wave.open(str(wav_path), 'w') as wf:
        wf.setnchannels(1)
        wf.setsampwidth(2)
        wf.setframerate(framerate)
        for i in range(int(framerate * duration)):
            val = int(1000 * math.sin(2 * math.pi * 440 * i / framerate))
            wf.writeframesraw(struct.pack('<h', val))
    (base / 'oral.txt').write_text('In this oral summary, the student explains the observed accuracy improvement and the main limitation.', encoding='utf-8')

    # mp4 video with transcript
    import cv2
    import numpy as np
    video_path = str(base / 'demo.mp4')
    fourcc = cv2.VideoWriter_fourcc(*'mp4v')
    writer = cv2.VideoWriter(video_path, fourcc, 2.0, (160, 120))
    for k in range(4):
        frame = np.full((120, 160, 3), 255 - 30 * k, dtype=np.uint8)
        cv2.putText(frame, f'Frame {k}', (20, 60), cv2.FONT_HERSHEY_SIMPLEX, 0.6, (0, 0, 0), 2)
        writer.write(frame)
    writer.release()
    (base / 'demo.txt').write_text('The video demonstration shows a simple four-frame walkthrough of the result.', encoding='utf-8')

    # xlsx dataset
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(['feature1', 'feature2', 'label'])
    ws.append([1.2, 0.5, 'A'])
    ws.append([1.5, 0.7, 'B'])
    ws.append([1.1, 0.4, 'A'])
    wb.save(str(base / 'dataset.xlsx'))

    # json dataset
    import json as _json
    (base / 'records.json').write_text(_json.dumps([{'x': 1, 'y': 2, 'target': 0}, {'x': 3, 'y': 4, 'target': 1}]), encoding='utf-8')


def build_mathematics_handwritten_like_submission(root: str) -> None:
    base = Path(root)
    base.mkdir(parents=True, exist_ok=True)
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (1200, 700), "white")
    draw = ImageDraw.Draw(img)
    lines = [
        "Proof: sum of two even integers is even",
        "Let a = 2m and b = 2n",
        "Then a + b = 2m + 2n = 2(m+n)",
        "Since m+n is an integer, a+b is even",
        "Hence proved",
    ]
    y = 40
    for line in lines:
        draw.text((40, y), line, fill="black")
        y += 110
    img.save(base / "solution.png")
